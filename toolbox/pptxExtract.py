"""
pptx_to_markdown_txt.py

Extracts readable text from a PowerPoint presentation (.pptx) and saves it to a
Markdown-formatted .txt file. Intended for feeding slide content as context to LLMs.

What it extracts:
- Slide titles as "## Slide N: <Title>"
- Text box and placeholder content as bullet lists with indentation
- Tables as Markdown tables (first row treated as header by default)
- Speaker notes (optional, on by default), shown under "Notes:" as a blockquote

Features:
- Handles grouped shapes (recursively)
- Preserves basic inline formatting (bold/italic) and hyperlinks where available
- Skips common non-content placeholders (slide number, date, footer)
- Escapes Markdown special characters so content is stable
- Replaces line breaks within table cells with <br> to keep tables valid

Usage:
  python pptx_to_markdown_txt.py slides.pptx [-o output.txt]
Options:
  --no-notes           Do not include speaker notes.
  --no-table-header    Do not treat first row of tables as a header.
  --max-slides N       Limit number of slides to export (0 = no limit).
  --encoding ENC       Output encoding (default: utf-8).

Requirements:
  pip install python-pptx
"""

import argparse
import os
import re
import sys

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
except ImportError as e:
    print("Missing dependency: python-pptx. Install it with:\n  pip install python-pptx", file=sys.stderr)
    sys.exit(1)


def md_escape(text: str, escape_pipes: bool = False) -> str:
    """
    Escape Markdown special characters in text.
    """
    if not text:
        return ""
    # Escape backslash first to avoid double-escaping.
    text = text.replace("\\", "\\\\")
    # Characters to escape in MD (avoid '-' to not disturb hyphenated words)
    chars = r"`*_{}[]()#+.!>"
    if escape_pipes:
        chars += r"|"
    return re.sub(r"([%s])" % re.escape(chars), r"\\\1", text)


def run_to_markdown(run) -> str:
    """
    Convert a pptx.text.text._Run to Markdown with basic inline formatting and hyperlinks.
    """
    t = run.text or ""
    if not t:
        return ""
    esc = md_escape(t)
    # Inline styles
    bold = getattr(run.font, "bold", None) is True
    italic = getattr(run.font, "italic", None) is True
    if bold and italic:
        esc = f"***{esc}***"
    elif bold:
        esc = f"**{esc}**"
    elif italic:
        esc = f"*{esc}*"

    # Hyperlink (if present on the run)
    url = None
    try:
        hl = getattr(run, "hyperlink", None)
        if hl:
            url = getattr(hl, "address", None) or getattr(hl, "target", None)
    except Exception:
        url = None

    if url:
        return f"[{esc}]({url})"
    return esc


def paragraph_to_markdown_text(paragraph) -> str:
    """
    Convert a pptx Paragraph (with runs) into Markdown inline string.
    Replaces internal newlines with a space for a single-line output per paragraph.
    """
    parts = []
    for r in paragraph.runs:
        parts.append(run_to_markdown(r))
    s = "".join(parts)
    # Normalize whitespace within a paragraph
    s = s.replace("\r\n", "\n").replace("\r", "\n")
    # Keep paragraph as a single line
    s = re.sub(r"\s*\n\s*", " / ", s)  # indicate explicit line breaks minimally
    return s.strip()


def text_frame_to_bullets(text_frame) -> list[str]:
    """
    Convert a text frame into a list of Markdown bullet lines, respecting indent level.
    """
    lines: list[str] = []
    for p in text_frame.paragraphs:
        text = paragraph_to_markdown_text(p)
        if not text:
            continue
        level = getattr(p, "level", 0) or 0
        indent = "  " * int(level)
        lines.append(f"{indent}- {text}")
    return lines


def escape_md_table_cell(text: str) -> str:
    """
    Escape/normalize text so it is safe inside a Markdown pipe table cell.
    """
    if text is None:
        return ""
    s = str(text)
    s = s.replace("\r\n", "\n").replace("\r", "\n")
    s = s.replace("\t", "    ")
    s = s.replace("\u00A0", " ")
    # Replace line breaks with <br> so the table stays valid
    s = s.replace("\n", "<br>")
    # Escape vertical bars so they don't split the table
    s = s.replace("|", r"\|")
    # Escape other MD specials lightly
    s = md_escape(s, escape_pipes=True)
    return s


def table_to_markdown(table, header=True) -> str:
    """
    Convert a pptx table to Markdown pipe table. First row is header by default.
    """
    rows: list[list[str]] = []
    for r in table.rows:
        row_cells: list[str] = []
        for c in r.cells:
            # Each cell may have multiple paragraphs
            para_texts = [paragraph_to_markdown_text(p) for p in c.text_frame.paragraphs]
            cell_text = " <br> ".join([t for t in para_texts if t])
            row_cells.append(escape_md_table_cell(cell_text))
        rows.append(row_cells)

    if not rows:
        return ""

    num_cols = max((len(r) for r in rows), default=0)
    # Normalize
    norm = [r + [""] * (num_cols - len(r)) for r in rows]

    if header and norm:
        hdr = norm[0]
        body = norm[1:]
        if not any(x.strip() for x in hdr):
            # If first row is empty, synthesize headers
            hdr = [f"Column {i+1}" for i in range(num_cols)]
            body = norm
    else:
        hdr = [f"Column {i+1}" for i in range(num_cols)]
        body = norm

    sep = ["---"] * num_cols
    lines = []
    lines.append("| " + " | ".join(hdr) + " |")
    lines.append("| " + " | ".join(sep) + " |")
    for r in body:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


SKIP_PLACEHOLDERS = {
    getattr(PP_PLACEHOLDER, "SLIDE_NUMBER", None),
    getattr(PP_PLACEHOLDER, "DATETIME", None),
    getattr(PP_PLACEHOLDER, "FOOTER", None),
    getattr(PP_PLACEHOLDER, "HEADER", None),
}


def is_skippable_placeholder(shape) -> bool:
    """
    Skip non-content placeholders like slide number, date, footer.
    """
    try:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            return ph and ph.type in SKIP_PLACEHOLDERS
    except Exception:
        pass
    return False


def iter_all_shapes(shapes):
    """
    Yield all shapes, descending into groups.
    """
    for sh in shapes:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub in iter_all_shapes(sh.shapes):
                    yield sub
            else:
                yield sh
        except Exception:
            # In case shape access fails, skip it
            continue


def get_slide_title(slide) -> str:
    """
    Get slide title text if available, else empty string.
    """
    try:
        title_shape = slide.shapes.title
    except Exception:
        title_shape = None
    if title_shape is not None:
        try:
            t = title_shape.text or ""
            t = t.strip()
            # Collapse internal whitespace
            t = re.sub(r"\s+", " ", t)
            return t
        except Exception:
            return ""
    # Fallback: first placeholder that looks like a title
    for sh in slide.shapes:
        try:
            if sh.is_placeholder and getattr(sh.placeholder_format, "type", None) in (
                getattr(PP_PLACEHOLDER, "TITLE", None),
                getattr(PP_PLACEHOLDER, "CENTER_TITLE", None),
            ):
                t = (sh.text or "").strip()
                t = re.sub(r"\s+", " ", t)
                if t:
                    return t
        except Exception:
            continue
    return ""


def extract_slide_content(slide, title_shape_id=None, include_tables=True, table_header=True):
    """
    Extract text bullets and tables from a slide (excluding title shape).
    Returns tuple (bullet_lines: list[str], tables_md: list[str], subtitles: list[str])
    """
    bullet_lines: list[str] = []
    tables_md: list[str] = []
    subtitles: list[str] = []

    for sh in iter_all_shapes(slide.shapes):
        # Skip the title shape if identified
        try:
            if title_shape_id is not None and getattr(sh, "shape_id", None) == title_shape_id:
                continue
        except Exception:
            pass

        # Skip common non-content placeholders
        if is_skippable_placeholder(sh):
            continue

        # Subtitle placeholder captured separately (as a line)
        try:
            if getattr(sh, "is_placeholder", False):
                ph_type = getattr(getattr(sh, "placeholder_format", None), "type", None)
                if ph_type == getattr(PP_PLACEHOLDER, "SUBTITLE", None):
                    if getattr(sh, "has_text_frame", False):
                        tlines = [paragraph_to_markdown_text(p) for p in sh.text_frame.paragraphs]
                        t = " ".join([x for x in tlines if x]).strip()
                        if t:
                            subtitles.append(t)
                    continue
        except Exception:
            pass

        # Tables
        try:
            if include_tables and getattr(sh, "has_table", False):
                md = table_to_markdown(sh.table, header=table_header)
                if md.strip():
                    tables_md.append(md)
                continue
        except Exception:
            pass

        # Text frames
        try:
            if getattr(sh, "has_text_frame", False):
                lines = text_frame_to_bullets(sh.text_frame)
                bullet_lines.extend(lines)
        except Exception:
            pass

    return bullet_lines, tables_md, subtitles


def extract_notes(slide) -> str | None:
    """
    Extract speaker notes as a single multi-line string. Returns None if absent.
    """
    ns = None
    try:
        ns = slide.notes_slide
    except Exception:
        ns = None
    if ns is None:
        return None

    # Preferred: dedicated notes_text_frame if present
    try:
        ntf = getattr(ns, "notes_text_frame", None)
        if ntf is not None:
            txt = (ntf.text or "").strip()
            if txt:
                return txt
    except Exception:
        pass

    # Fallback: aggregate from shapes that have text
    parts = []
    try:
        for sh in iter_all_shapes(ns.shapes):
            if is_skippable_placeholder(sh):
                continue
            if getattr(sh, "has_text_frame", False):
                for p in sh.text_frame.paragraphs:
                    t = paragraph_to_markdown_text(p)
                    if t:
                        parts.append(t)
    except Exception:
        pass

    text = "\n".join(parts).strip()
    return text if text else None


def main(inputPath):
    ap = argparse.ArgumentParser(description="Extract text from a PowerPoint (.pptx) into Markdown-formatted .txt for LLM context.")
    ap.add_argument("--no-notes", action="store_true", help="Do not include speaker notes.")
    ap.add_argument("--no-table-header", action="store_true", help="Do not treat first table row as a header.")
    ap.add_argument("--max-slides", type=int, default=0, help="Limit number of slides to export (0 = no limit).")
    ap.add_argument("--encoding", default="utf-8", help="Output encoding (default: utf-8).")
    args = ap.parse_args()

    input_path = inputPath
    if not os.path.isfile(input_path):
        print(f"Error: File not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    _, ext = os.path.splitext(input_path)
    if ext.lower() != ".pptx":
        print("Error: Only .pptx files are supported.", file=sys.stderr)
        sys.exit(1)

    try:
        prs = Presentation(input_path)
    except Exception as e:
        print(f"Failed to open presentation: {e}", file=sys.stderr)
        sys.exit(1)

    out_lines: list[str] = []
    out_lines.append("")

    max_slides = args.max_slides if args.max_slides and args.max_slides > 0 else None

    for idx, slide in enumerate(prs.slides, start=1):
        if max_slides is not None and idx > max_slides:
            break

        # Title
        title_text = get_slide_title(slide)
        title_shape_id = None
        try:
            title_shape = slide.shapes.title
            if title_shape is not None:
                title_shape_id = getattr(title_shape, "shape_id", None)
        except Exception:
            pass

        header = f"## Slide {idx}"
        if title_text:
            header += f": {md_escape(title_text)}"
        out_lines.append(header)
        out_lines.append("")

        # Content extraction
        bullets, tables, subtitles = extract_slide_content(
            slide,
            title_shape_id=title_shape_id,
            include_tables=True,
            table_header=(not args.no_table_header),
        )

        # Subtitle(s)
        for sub in subtitles:
            out_lines.append(f"_Subtitle:_ {md_escape(sub)}")
        if subtitles:
            out_lines.append("")

        # Bullet lines (text)
        if bullets:
            out_lines.extend(bullets)
            out_lines.append("")

        # Tables
        for ti, md in enumerate(tables, start=1):
            out_lines.append(f"Table {ti}:")
            out_lines.append(md)
            out_lines.append("")

        # Notes
        if not args.no_notes:
            notes = extract_notes(slide)
            if notes:
                out_lines.append("Notes:")
                for line in notes.splitlines():
                    line = line.strip()
                    if not line:
                        continue
                    out_lines.append(f"> {md_escape(line)}")
                out_lines.append("")

    content = "\n".join(out_lines).rstrip() + "\n"

    return content

